import os
import shutil
import tempfile
from unittest.mock import MagicMock, Mock, mock_open, patch

import pytest
import requests
from PIL import Image
from pptx import Presentation

from src.file_parser.other_files_parser import FileConverter


@patch('src.file_parser.other_files_parser.logging.getLogger')
class TestFileConverter:
    """Test suite for FileConverter class"""

    @pytest.fixture
    def temp_dir(self):
        """Create a temporary directory for test files"""
        temp_dir = tempfile.mkdtemp()
        yield temp_dir
        shutil.rmtree(temp_dir, ignore_errors=True)

    @pytest.fixture
    def sample_files(self, temp_dir):
        """Create sample test files"""
        files = {}

        img_path = os.path.join(temp_dir, 'test_image.jpg')
        img = Image.new('RGB', (100, 100), color='red')
        img.save(img_path, 'JPEG')
        files['image'] = img_path

        html_path = os.path.join(temp_dir, 'test.html')
        with open(html_path, 'w') as f:
            f.write('<html><body><h1>Test HTML</h1></body></html>')
        files['html'] = html_path

        md_path = os.path.join(temp_dir, 'test.md')
        with open(md_path, 'w') as f:
            f.write('# Test Markdown\n\nThis is a test.')
        files['markdown'] = md_path

        txt_path = os.path.join(temp_dir, 'test.txt')
        with open(txt_path, 'w') as f:
            f.write('This is a text file.')
        files['text'] = txt_path

        pptx_path = os.path.join(temp_dir, 'test.pptx')

        prs = Presentation()
        prs.save(pptx_path)
        files['pptx'] = pptx_path

        return files

    def test_init(self, mock_get_logger, temp_dir):
        """Test FileConverter initialization"""
        mock_logger = MagicMock()
        mock_get_logger.return_value = mock_logger
        file_path = os.path.join(temp_dir, 'test.pdf')
        converter = FileConverter(file_path, temp_dir)

        assert converter.file_path == file_path
        assert converter.output_dir == temp_dir
        assert converter.temp_files == []
        assert os.path.exists(temp_dir)

    def test_init_default_output_dir(self, mock_get_logger, temp_dir):
        """Test FileConverter initialization with default output directory"""
        file_path = os.path.join(temp_dir, 'test.pdf')
        with patch('os.makedirs') as mock_makedirs:
            converter = FileConverter(file_path)
            assert converter.output_dir == 'assets'
            mock_makedirs.assert_called_once_with('assets', exist_ok=True)

    def test_detect_file_type_nonexistent_file(self, mock_get_logger):
        """Test file type detection for non-existent file"""
        converter = FileConverter('/nonexistent/file.pdf')

        with pytest.raises(FileNotFoundError):
            converter.detect_file_type()

    def test_detect_file_type_by_extension(
        self,
        mock_get_logger,
        sample_files,
        temp_dir,
    ):
        """Test file type detection by extension"""
        test_cases = [
            (sample_files['image'], 'images'),
            (sample_files['html'], 'web'),
            (sample_files['markdown'], 'markdown'),
        ]

        for file_path, expected_type in test_cases:
            converter = FileConverter(file_path, temp_dir)
            assert converter.detect_file_type() == expected_type

    @patch('mimetypes.guess_type')
    def test_detect_file_type_by_mime(
        self,
        mock_guess_type,
        mock_get_logger,
        sample_files,
        temp_dir,
    ):
        """Test file type detection by MIME type"""

        mock_guess_type.return_value = ('image/jpeg', None)
        converter = FileConverter(
            sample_files['text'],
            temp_dir,
        )  # Use .txt file
        assert converter.detect_file_type() == 'images'

        mock_guess_type.return_value = ('application/pdf', None)
        assert converter.detect_file_type() == 'pdf'

    def test_detect_file_type_unknown(
        self,
        mock_get_logger,
        sample_files,
        temp_dir,
    ):
        """Test file type detection for unknown file"""
        converter = FileConverter(sample_files['text'], temp_dir)
        assert converter.detect_file_type() == 'unknown'

    def test_generate_unique_filename(self, mock_get_logger, temp_dir):
        """Test unique filename generation"""
        converter = FileConverter('test.pdf', temp_dir)

        filename1 = converter._generate_unique_filename('test', '.pdf')
        expected_path1 = os.path.join(temp_dir, 'test.pdf')
        assert filename1 == expected_path1

        with open(expected_path1, 'w') as f:
            f.write('test')

        with patch('time.time', return_value=1234567890):
            filename2 = converter._generate_unique_filename('test', '.pdf')
            expected_path2 = os.path.join(temp_dir, 'test_1234567890.pdf')
            assert filename2 == expected_path2

    @patch('PIL.Image.open')
    def test_convert_image_to_pdf_success(
        self,
        mock_image_open,
        mock_get_logger,
        temp_dir,
    ):
        """Test successful image to PDF conversion"""
        mock_img = Mock()
        mock_img.mode = 'RGB'
        mock_img.__enter__ = Mock(return_value=mock_img)
        mock_img.__exit__ = Mock(return_value=None)
        mock_image_open.return_value = mock_img

        converter = FileConverter('test.jpg', temp_dir)

        with patch.object(
            converter,
            '_generate_unique_filename',
        ) as mock_filename:
            output_path = os.path.join(temp_dir, 'test.pdf')
            mock_filename.return_value = output_path

            result = converter.convert_image_to_pdf()

            assert result == output_path
            mock_img.save.assert_called_once_with(
                output_path,
                'PDF',
                resolution=100.0,
            )

    @patch('PIL.Image.open')
    def test_convert_image_to_pdf_mode_conversion(
        self,
        mock_image_open,
        mock_get_logger,
        temp_dir,
    ):
        """Test image to PDF conversion with mode conversion"""
        mock_img = Mock()
        mock_img.mode = 'RGBA'
        mock_converted_img = Mock()
        mock_img.convert.return_value = mock_converted_img
        mock_img.__enter__ = Mock(return_value=mock_img)
        mock_img.__exit__ = Mock(return_value=None)
        mock_image_open.return_value = mock_img

        converter = FileConverter('test.png', temp_dir)

        with patch.object(
            converter,
            '_generate_unique_filename',
        ) as mock_filename:
            output_path = os.path.join(temp_dir, 'test.pdf')
            mock_filename.return_value = output_path

            converter.convert_image_to_pdf()

            mock_img.convert.assert_called_once_with('RGB')
            mock_converted_img.save.assert_called_once()

    @patch('PIL.Image.open')
    def test_convert_image_to_pdf_failure(
        self,
        mock_image_open,
        mock_get_logger,
        temp_dir,
    ):
        """Test image to PDF conversion failure"""
        mock_image_open.side_effect = Exception('Image open failed')

        converter = FileConverter('test.jpg', temp_dir)

        with pytest.raises(Exception, match='Image conversion error'):
            converter.convert_image_to_pdf()

    def test_is_valid_url(self, mock_get_logger, temp_dir):
        """Test URL validation"""
        converter = FileConverter('test.pdf', temp_dir)

        assert converter.is_valid_url('https://www.example.com') is True
        assert converter.is_valid_url('http://example.com') is True
        assert (
            converter.is_valid_url(
                'https://subdomain.example.com/path',
            )
            is True
        )

        assert converter.is_valid_url('not-a-url') is False
        assert converter.is_valid_url('') is False
        assert converter.is_valid_url('example.com') is False

    def test_get_domain_name(self, mock_get_logger, temp_dir):
        """Test domain name extraction"""
        converter = FileConverter('test.pdf', temp_dir)

        test_cases = [
            ('https://www.example.com', 'example_com'),
            ('https://subdomain.example.com/path', 'subdomain_example_com'),
            ('http://test-site.co.uk', 'test_site_co_uk'),
        ]

        for url, expected in test_cases:
            assert converter.get_domain_name(url) == expected

    def test_get_domain_name_error(self, mock_get_logger, temp_dir):
        """Test domain name extraction with error"""
        converter = FileConverter('test.pdf', temp_dir)

        with patch(
            'src.file_parser.other_files_parser.urlparse',
            side_effect=Exception('Parse error'),
        ):
            result = converter.get_domain_name('invalid-url')
            assert result == 'website'

    @patch('requests.get')
    @patch('PyQt5.QtWidgets.QApplication')
    def test_convert_url_to_pdf_invalid_url(
        self,
        mock_qapp,
        mock_requests,
        mock_get_logger,
        temp_dir,
    ):
        """Test URL to PDF conversion with invalid URL"""
        converter = FileConverter('not-a-url', temp_dir)

        with pytest.raises(ValueError, match='Invalid URL'):
            converter.convert_url_to_pdf()

    @patch('requests.get')
    def test_convert_url_to_pdf_request_error(
        self,
        mock_requests,
        mock_get_logger,
        temp_dir,
    ):
        """Test URL to PDF conversion with request error"""
        mock_requests.side_effect = requests.RequestException(
            'Connection failed',
        )
        converter = FileConverter('https://example.com', temp_dir)

        with pytest.raises(Exception, match='Connection error'):
            converter.convert_url_to_pdf()

    @patch('requests.get')
    def test_convert_url_to_pdf_bad_status(
        self,
        mock_requests,
        mock_get_logger,
        temp_dir,
    ):
        """Test URL to PDF conversion with bad HTTP status"""
        mock_response = Mock()
        mock_response.status_code = 404
        mock_requests.return_value = mock_response

        converter = FileConverter('https://example.com', temp_dir)

        with pytest.raises(Exception, match='Unavailable page'):
            converter.convert_url_to_pdf()

    @patch('pdfkit.from_file')
    def test_convert_html_to_pdf_success(
        self,
        mock_pdfkit,
        mock_get_logger,
        temp_dir,
    ):
        """Test successful HTML to PDF conversion"""
        converter = FileConverter('test.html', temp_dir)

        with patch.object(
            converter,
            '_generate_unique_filename',
        ) as mock_filename:
            output_path = os.path.join(temp_dir, 'test.pdf')
            mock_filename.return_value = output_path

            result = converter.convert_html_to_pdf()

            assert result == output_path
            mock_pdfkit.assert_called_once_with(
                'test.html',
                output_path,
                options={'quiet': ''},
            )

    @patch('pdfkit.from_file')
    def test_convert_html_to_pdf_failure(
        self,
        mock_pdfkit,
        mock_get_logger,
        temp_dir,
    ):
        """Test HTML to PDF conversion failure"""
        mock_pdfkit.side_effect = Exception('PDFKit failed')

        converter = FileConverter('test.html', temp_dir)

        with pytest.raises(Exception, match='HTML conversion error'):
            converter.convert_html_to_pdf()

    @patch('pdfkit.from_string')
    @patch('markdown.markdown')
    def test_convert_markdown_to_pdf_success(
        self,
        mock_markdown,
        mock_pdfkit,
        mock_get_logger,
        temp_dir,
    ):
        """Test successful Markdown to PDF conversion"""
        mock_markdown.return_value = '<h1>Test</h1>'

        converter = FileConverter('test.md', temp_dir)

        with patch(
            'builtins.open',
            mock_open(read_data='# Test'),
        ) as mock_file:
            with patch.object(
                converter,
                '_generate_unique_filename',
            ) as mock_filename:
                output_path = os.path.join(temp_dir, 'test.pdf')
                mock_filename.return_value = output_path

                result = converter.convert_markdown_to_pdf()

                assert result == output_path
                mock_file.assert_called_once_with('test.md', encoding='utf-8')
                mock_markdown.assert_called_once_with('# Test')
                mock_pdfkit.assert_called_once_with(
                    '<h1>Test</h1>',
                    output_path,
                    options={'quiet': ''},
                )

    @patch('pdfkit.from_string')
    @patch('markdown.markdown')
    def test_convert_markdown_to_pdf_failure(
        self,
        mock_markdown,
        mock_pdfkit,
        mock_get_logger,
        temp_dir,
    ):
        """Test Markdown to PDF conversion failure"""
        mock_pdfkit.side_effect = Exception('PDFKit failed')

        converter = FileConverter('test.md', temp_dir)

        with patch('builtins.open', mock_open(read_data='# Test')):
            with pytest.raises(Exception, match='Markdown conversion error'):
                converter.convert_markdown_to_pdf()

    @patch('reportlab.pdfgen.canvas.Canvas')
    @patch('pptx.Presentation')
    def test_convert_pptx_to_pdf_success(
        self,
        mock_presentation,
        mock_canvas,
        mock_get_logger,
        temp_dir,
        sample_files,
    ):
        """Test successful PPTX to PDF conversion"""

        mock_slide1 = Mock()
        mock_slide2 = Mock()
        mock_shape = Mock()
        mock_shape.text = 'Test text\nSecond line'
        mock_slide1.shapes = [mock_shape]
        mock_slide2.shapes = []

        mock_prs = Mock()
        mock_prs.slides = [mock_slide1, mock_slide2]
        mock_presentation.return_value = mock_prs

        mock_canvas_instance = Mock()
        mock_canvas.return_value = mock_canvas_instance

        with patch('src.file_parser.other_files_parser.letter', (612, 792)):
            converter = FileConverter(sample_files['pptx'], temp_dir)

            with patch.object(
                converter,
                '_generate_unique_filename',
            ) as mock_filename:
                output_path = os.path.join(temp_dir, 'test.pdf')
                mock_filename.return_value = output_path

                result = converter.convert_pptx_to_pdf()

                assert result == output_path
                mock_canvas_instance.save.assert_called_once()

    @patch('pptx.Presentation')
    def test_convert_pptx_to_pdf_failure(
        self,
        mock_presentation,
        mock_get_logger,
        temp_dir,
    ):
        """Test PPTX to PDF conversion failure"""
        mock_presentation.side_effect = Exception('PPTX failed')

        converter = FileConverter('test.pptx', temp_dir)

        with pytest.raises(Exception, match='PPTX conversion error'):
            converter.convert_pptx_to_pdf()

    def test_convert_to_pdf_url(self, mock_get_logger, temp_dir):
        """Test convert_to_pdf with URL"""
        converter = FileConverter('https://example.com', temp_dir)

        with patch.object(converter, 'convert_url_to_pdf') as mock_convert:
            mock_convert.return_value = 'output.pdf'
            result = converter.convert_to_pdf()
            assert result == 'output.pdf'
            mock_convert.assert_called_once()

    def test_convert_to_pdf_already_pdf(self, mock_get_logger, temp_dir):
        """Test convert_to_pdf with PDF file"""
        pdf_path = os.path.join(temp_dir, 'test.pdf')
        with open(pdf_path, 'w') as f:
            f.write('dummy pdf')

        converter = FileConverter(pdf_path, temp_dir)

        with patch.object(converter, 'detect_file_type', return_value='pdf'):
            result = converter.convert_to_pdf()
            assert result == pdf_path

    def test_convert_to_pdf_image(self, mock_get_logger, temp_dir):
        """Test convert_to_pdf with image file"""
        converter = FileConverter('test.jpg', temp_dir)

        with patch.object(
            converter,
            'detect_file_type',
            return_value='images',
        ):
            with patch.object(
                converter,
                'convert_image_to_pdf',
            ) as mock_convert:
                mock_convert.return_value = 'output.pdf'
                result = converter.convert_to_pdf()
                assert result == 'output.pdf'
                mock_convert.assert_called_once()

    def test_convert_to_pdf_unsupported(self, mock_get_logger, temp_dir):
        """Test convert_to_pdf with unsupported file type"""
        converter = FileConverter('test.txt', temp_dir)

        with patch.object(
            converter,
            'detect_file_type',
            return_value='unknown',
        ):
            with pytest.raises(ValueError, match='Unsupported file type'):
                converter.convert_to_pdf()

    def test_cleanup(self, mock_get_logger, temp_dir):
        """Test cleanup method"""
        converter = FileConverter('test.pdf', temp_dir)

        temp_file1 = os.path.join(temp_dir, 'temp1.txt')
        temp_file2 = os.path.join(temp_dir, 'temp2.txt')

        with open(temp_file1, 'w') as f:
            f.write('temp1')
        with open(temp_file2, 'w') as f:
            f.write('temp2')

        converter.temp_files = [temp_file1, temp_file2]

        converter.cleanup()

        assert not os.path.exists(temp_file1)
        assert not os.path.exists(temp_file2)
        assert converter.temp_files == []

    def test_cleanup_with_missing_file(self, mock_get_logger, temp_dir):
        """Test cleanup with missing temporary file"""
        converter = FileConverter('test.pdf', temp_dir)

        non_existent_file = os.path.join(temp_dir, 'non_existent.txt')
        converter.temp_files = [non_existent_file]

        converter.cleanup()
        assert converter.temp_files == []

    @patch('src.file_parser.other_files_parser.pdf_parser.PdfParser')
    def test_initiate_parser_success(
        self,
        mock_pdf_parser,
        mock_get_logger,
        temp_dir,
        sample_files,
    ):
        """Test successful parser initiation"""
        mock_parser_instance = mock_pdf_parser.return_value
        mock_parser_instance.initiate.return_value = 'Parsed LLM Content'

        converter = FileConverter(sample_files['image'], temp_dir)
        with patch.object(
            converter,
            'convert_to_pdf',
            return_value='converted.pdf',
        ) as mock_convert:
            result = converter.initiate_parser()

            assert result == 'Parsed LLM Content'
            mock_convert.assert_called_once()
            mock_pdf_parser.assert_called_once_with(
                file_path='converted.pdf',
                output_dir='extracted_content',
                describe_images=True,
            )
            mock_parser_instance.initiate.assert_called_once()

    def test_initiate_parser_failure(
        self,
        mock_get_logger,
        temp_dir,
        sample_files,
    ):
        """Test parser initiation failure"""
        with patch(
            'src.file_parser.other_files_parser.FileConverter.convert_to_pdf',
            side_effect=Exception('Conversion failed'),
        ):
            converter = FileConverter(sample_files['image'], temp_dir)
            with pytest.raises(Exception, match='Failed to parse file'):
                converter.initiate_parser()

    def test_del_calls_cleanup(self, mock_get_logger, temp_dir):
        """Test that __del__ calls cleanup"""
        with patch.object(FileConverter, 'cleanup') as mock_cleanup:
            converter = FileConverter('test.pdf', temp_dir)
            converter.__del__()
            mock_cleanup.assert_called_once()
