"""
This module provides the FileConverter class for converting various file types (images, HTML, Markdown, PPTX, URLs) to PDF and initiating PDF parsing workflows.
"""
from __future__ import annotations

import gc
import logging
import mimetypes
import os
import subprocess
import sys
import time
from urllib.parse import urlparse

import markdown
import pdfkit
import requests
from PIL import Image
from pptx import Presentation
from reportlab.lib.pagesizes import letter
from reportlab.pdfgen import canvas

from src.file_parser import pdf_parser
from src.utils.logging_config import get_request_id, get_session_logger


class FileConverter:
    """
    Handles conversion of supported file types (images, HTML, Markdown, PPTX, URLs) to PDF format and initiates PDF parsing.
    """

    SUPPORTED_FORMATS = {
        "images": [".jpg", ".jpeg", ".png", ".bmp", ".tiff", ".gif"],
        "web": [".html", ".htm"],
        "presentations": [".pptx"],
        "markdown": [".md", ".markdown"],
        "pdf": [".pdf"],
    }

    def __init__(self, file_path: str, output_dir: str = 'assets', request_id: str | None = None):
        """
        Initialize the FileConverter.

        Args:
            file_path (str): Path to the input file or URL.
            output_dir (str, optional): Directory to store output files. Defaults to 'assets'.
            request_id (str, optional): Unique request identifier for logging. If None, a new one is generated.
        """
        self.request_id = request_id or get_request_id()
        self.logger = get_session_logger(self.request_id)
        self.file_path = file_path
        self.output_dir = output_dir
        self.temp_files: list[str] = []
        self.create_output_dir()

    def create_output_dir(self) -> None:
        """
        Create the output directory if it doesn't exist.

        Raises:
            OSError: If the directory cannot be created.
        """
        try:
            os.makedirs(self.output_dir, exist_ok=True)
        except OSError as e:
            self.logger.error(
                f"Failed to create output directory {self.output_dir}: {e}",
            )
            raise

    def detect_file_type(self) -> str:
        """
        Detect the type of file based on its extension or MIME type.

        Returns:
            str: Detected file type (e.g., 'images', 'web', 'presentations', 'markdown', 'pdf', or 'unknown').

        Raises:
            FileNotFoundError: If the file does not exist.
        """
        if not os.path.exists(self.file_path):
            self.logger.error(f"File does not exist: {self.file_path}")
            raise FileNotFoundError(f"File does not exist: {self.file_path}")

        file_ext = os.path.splitext(self.file_path)[1].lower()
        self.logger.debug(f"Detected file extension: {file_ext}")

        for file_type, extensions in self.SUPPORTED_FORMATS.items():
            if file_ext in extensions:
                self.logger.info(f"File type detected: {file_type}")
                return file_type

        mime_type, _ = mimetypes.guess_type(self.file_path)
        if mime_type:
            self.logger.debug(f"MIME type detected: {mime_type}")
            if mime_type.startswith("image/"):
                self.logger.info("File type detected: images (via MIME)")
                return "images"
            elif mime_type == "application/pdf":
                self.logger.info("File type detected: pdf (via MIME)")
                return "pdf"

        self.logger.warning(f"Unknown file type for: {self.file_path}")
        return "unknown"

    def _generate_unique_filename(
        self,
        base_name: str,
        extension: str,
    ) -> str:
        """
        Generate a unique filename in the output directory.

        Args:
            base_name (str): Base name for the file.
            extension (str): File extension (e.g., '.pdf').

        Returns:
            str: Full path to the unique output file.
        """
        output_path = os.path.join(self.output_dir, f'{base_name}{extension}')

        if os.path.exists(output_path):
            timestamp = int(time.time())
            output_path = os.path.join(
                self.output_dir,
                f"{base_name}_{timestamp}{extension}",
            )
            self.logger.debug(f"Generated unique filename: {output_path}")

        return output_path

    def convert_image_to_pdf(self) -> str:
        """
        Convert an image file to PDF format.

        Returns:
            str: Path to the generated PDF file.

        Raises:
            RuntimeError: If image conversion fails.
        """
        self.logger.info(f'Converting image to PDF: {self.file_path}')
        original_name = os.path.splitext(os.path.basename(self.file_path))[0]
        output_path = self._generate_unique_filename(original_name, ".pdf")

        try:
            with Image.open(self.file_path) as img:
                if img.mode != "RGB":
                    img = img.convert("RGB")
                    self.logger.debug("Image converted to RGB mode")
                img.save(output_path, "PDF", resolution=100.0)
            self.logger.info(
                f"Image successfully converted to PDF: {output_path}",
            )
            return output_path
        except (OSError, ValueError) as e:
            self.logger.error(f"Image conversion failed: {str(e)}")
            raise RuntimeError(f"Image conversion error: {str(e)}")

    def is_valid_url(self, url: str) -> bool:
        """
        Validate if the provided string is a valid URL.

        Args:
            url (str): URL string to validate.

        Returns:
            bool: True if valid, False otherwise.
        """
        try:
            result = urlparse(url)
            is_valid = all([result.scheme, result.netloc])
            self.logger.debug(f"URL validation for {url}: {is_valid}")
            return is_valid
        except ValueError as e:
            self.logger.debug(f"URL validation failed for {url}: {str(e)}")
            return False

    def get_domain_name(self, url: str) -> str:
        """
        Extract the domain name from a URL and format it for use in filenames.

        Args:
            url (str): The URL to extract the domain from.

        Returns:
            str: Cleaned domain name suitable for filenames.
        """
        try:
            domain = urlparse(url).netloc
            clean_domain = (
                domain.replace("www.", "")
                .replace(
                    ".",
                    "_",
                )
                .replace("-", "_")[:30]
            )
            self.logger.debug(f"Extracted domain name: {clean_domain}")
            return clean_domain
        except (AttributeError, ValueError) as e:
            self.logger.warning(
                f"Could not extract domain from {url}: {str(e)}",
            )
            return "website"

    def convert_url_to_pdf(self) -> str:
        """
        Convert a URL to a PDF file using an external script.

        Returns:
            str: Path to the generated PDF file.

        Raises:
            ValueError: If the URL is invalid.
            Exception: If the URL is inaccessible or conversion fails.
        """
        self.logger.info(f'Converting URL to PDF: {self.file_path}')

        if not self.is_valid_url(self.file_path):
            self.logger.error(f"Invalid URL: {self.file_path}")
            raise ValueError(f"Invalid URL: {self.file_path}")

        try:
            self.logger.debug("Checking URL accessibility...")
            response = requests.get(
                self.file_path,
                timeout=10,
                headers={
                    "User-Agent": "Mozilla/5.0",
                },
            )
            if response.status_code >= 400:
                self.logger.error(
                    f"URL returned status code {response.status_code}",
                )
                raise Exception(
                    f"Unavailable page (code: {response.status_code})",
                )
        except requests.RequestException as e:
            self.logger.error(
                f"Connection error for URL {self.file_path}: {str(e)}",
            )
            raise Exception(f"Connection error: {str(e)}")

        domain_name = self.get_domain_name(self.file_path)
        output_path = os.path.abspath(
            self._generate_unique_filename(domain_name, ".pdf")
        )

        os.makedirs(os.path.dirname(output_path), exist_ok=True)

        script_path = os.path.join(os.path.dirname(__file__), "url2pdf.py")
        try:
            subprocess.run(
                [sys.executable, script_path, self.file_path, output_path],
                check=True,
                capture_output=True,
                text=True,
            )
            # Now wait for the file to physically appear
            timeout = 10
            start_time = time.time()
            while not os.path.exists(output_path):
                if time.time() - start_time > timeout:
                    self.logger.error(
                        f"PDF conversion succeeded but file not found: {output_path}"
                    )
                    raise FileNotFoundError(
                        f"PDF not found after subprocess: {output_path}"
                    )
                time.sleep(0.1)

            self.logger.info(f"URL successfully converted to PDF: {output_path}")
            return output_path

        except subprocess.CalledProcessError as e:
            self.logger.error(f"URL to PDF conversion failed: {e.stderr}")
            raise Exception(f"URL to PDF conversion error: {e.stderr}")

    def convert_html_to_pdf(self) -> str:
        """
        Convert an HTML file to PDF format.

        Returns:
            str: Path to the generated PDF file.

        Raises:
            RuntimeError: If HTML conversion fails.
        """
        self.logger.info(f'Converting HTML to PDF: {self.file_path}')
        original_name = os.path.splitext(os.path.basename(self.file_path))[0]
        output_path = self._generate_unique_filename(original_name, ".pdf")

        try:
            pdfkit.from_file(
                self.file_path,
                output_path,
                options={"quiet": ""},
            )
            self.logger.info(
                f"HTML successfully converted to PDF: {output_path}",
            )
            return output_path
        except (OSError, ValueError) as e:
            self.logger.error(f"HTML conversion failed: {str(e)}")
            raise RuntimeError(f"HTML conversion error: {str(e)}")

    def convert_markdown_to_pdf(self) -> str:
        """
        Convert a Markdown file to PDF format.

        Returns:
            str: Path to the generated PDF file.

        Raises:
            RuntimeError: If Markdown conversion fails.
        """
        self.logger.info(f'Converting Markdown to PDF: {self.file_path}')
        original_name = os.path.splitext(os.path.basename(self.file_path))[0]
        output_path = self._generate_unique_filename(original_name, ".pdf")

        try:
            with open(self.file_path, encoding="utf-8") as f:
                html = markdown.markdown(f.read())
            pdfkit.from_string(html, output_path, options={"quiet": ""})
            self.logger.info(
                f"Markdown successfully converted to PDF: {output_path}",
            )
            return output_path
        except (OSError, ValueError) as e:
            self.logger.error(f"Markdown conversion failed: {str(e)}")
            raise RuntimeError(f"Markdown conversion error: {str(e)}")

    def convert_pptx_to_pdf(self) -> str:
        """
        Convert a PPTX file to PDF format.

        Returns:
            str: Path to the generated PDF file.

        Raises:
            RuntimeError: If PPTX conversion fails.
        """
        self.logger.info(f'Converting PPTX to PDF: {self.file_path}')
        original_name = os.path.splitext(os.path.basename(self.file_path))[0]
        output_path = self._generate_unique_filename(original_name, ".pdf")

        try:
            prs = Presentation(self.file_path)
            c = canvas.Canvas(output_path, pagesize=letter)
            _, height = letter

            self.logger.debug(f"Processing {len(prs.slides)} slides...")

            for i, slide in enumerate(prs.slides):
                if i > 0:
                    c.showPage()

                y = height - 40
                c.setFont("Helvetica", 12)
                c.drawString(40, y, f"Slide {i + 1}")
                y -= 30

                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        for line in shape.text.split("\n"):
                            if y < 40:
                                c.showPage()
                                y = height - 40
                            c.drawString(40, y, line[:100])
                            y -= 15
            c.save()
            self.logger.info(
                f"PPTX successfully converted to PDF: {output_path}",
            )
            return output_path
        except (OSError, ValueError) as e:
            self.logger.error(f"PPTX conversion failed: {str(e)}")
            raise RuntimeError(f"PPTX conversion error: {str(e)}")

    def convert_to_pdf(self) -> str:
        """
        Convert the input file (or URL) to PDF format, if needed.

        Returns:
            str: Path to the resulting PDF file.

        Raises:
            ValueError: If the file type is unsupported.
        """
        self.logger.info('Starting file conversion to PDF...')

        if self.is_valid_url(self.file_path):
            self.logger.info("Processing as URL")
            return self.convert_url_to_pdf()

        file_type = self.detect_file_type()

        if file_type == "pdf":
            self.logger.info("File is already PDF, no conversion needed")
            return self.file_path
        elif file_type == "images":
            return self.convert_image_to_pdf()
        elif file_type == "web":
            return self.convert_html_to_pdf()
        elif file_type == "markdown":
            return self.convert_markdown_to_pdf()
        elif file_type == "presentations":
            return self.convert_pptx_to_pdf()
        else:
            self.logger.error(f"Unsupported file type: {file_type}")
            raise ValueError(f"Unsupported file type: {file_type}")

    def cleanup(self):
        """
        Explicitly clean up all temporary files and resources used during conversion.
        """
        if self.temp_files:
            self.logger.debug(
                f"Cleaning up {len(self.temp_files)} temporary files...",
            )

        for file_path in self.temp_files[:]:
            try:
                if os.path.exists(file_path):
                    os.remove(file_path)
                    self.logger.debug(f"Removed temporary file: {file_path}")
                self.temp_files.remove(file_path)
            except OSError as e:
                self.logger.warning(
                    f"Could not remove temporary file {file_path}: {str(e)}",
                )

        gc.collect()
        self.logger.debug("Cleanup completed")

    def __del__(self):
        """
        Destructor to ensure cleanup of resources.
        """
        self.cleanup()

    def initiate_parser(self) -> str:
        """
        Initiate the parser workflow by converting the file to PDF and extracting its content using PdfParser.

        Returns:
            str: Content extracted from the PDF, formatted for LLMs.

        Raises:
            RuntimeError: If parsing or conversion fails.
        """
        self.logger.info("Initiating parser workflow...")

        try:
            pdf_path = self.convert_to_pdf()
            self.logger.info(f"File converted to PDF: {pdf_path}")

            self.logger.info("Initializing PDF parser...")
            pdf_parser_instance = pdf_parser.PdfParser(
                file_path=pdf_path,
                output_dir="extracted_content",
                describe_images=True,
                request_id=self.request_id,
            )

            self.logger.info("Extracting content from PDF...")
            llm_content = pdf_parser_instance.initiate()

            self.logger.info("Parser workflow completed successfully")
            return llm_content

        except (OSError, ValueError, RuntimeError) as e:
            self.logger.error(f"Parser workflow failed: {str(e)}")
            raise RuntimeError(
                f"Failed to parse file {self.file_path}: {str(e)}",
            )
        finally:
            self.cleanup()
